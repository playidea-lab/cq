"""PPTX Parser - python-pptx 기반."""

import logging
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

from c4.c2.parsers.base import BaseParser, ImageData, ParseResult
from c4.c2.parsers.ir_models import (
    Document,
    MergeInfo,
    create_heading,
    create_image,
    create_paragraph,
    create_table,
)
from c4.c2.parsers.utils.chart_parser import parse_chart_xml
from c4.c2.parsers.utils.image import generate_image_id, get_extension_from_mime

logger = logging.getLogger(__name__)


class PptxParser(BaseParser):
    """PPTX 문서 파서.

    파싱 규칙:
    - slide title → h1/h2
    - subtitle → h3
    - text box → p
    - table → header/body 분리 + 머지셀 지원
    - 슬라이드별 구분 (슬라이드 번호 표시)
    """

    @property
    def supported_extensions(self) -> list[str]:
        return [".pptx"]  # .ppt는 ppt_parser에서 처리

    def parse(self, file_path: Path) -> Document:
        """PPTX 파일을 IR로 변환 (이미지 제외)."""
        result = self.parse_with_images(file_path)
        return result.document

    def parse_with_images(self, file_path: Path) -> ParseResult:
        """PPTX 파일을 IR과 이미지로 변환."""
        prs = Presentation(str(file_path))
        blocks = []
        extracted_images: list[ImageData] = []
        image_index = 0

        for slide_idx, slide in enumerate(prs.slides):
            slide_blocks, slide_images, image_index = self._parse_slide_with_images(
                slide, slide_idx, image_index
            )
            blocks.extend(slide_blocks)
            extracted_images.extend(slide_images)

        # 차트 추출 (ZIP 구조에서)
        chart_blocks, chart_images = self._extract_charts(file_path, image_index)
        blocks.extend(chart_blocks)
        extracted_images.extend(chart_images)

        return ParseResult(document=Document(blocks=blocks), images=extracted_images)

    def _parse_slide_with_images(
        self, slide, slide_idx: int, image_index: int
    ) -> tuple[list, list[ImageData], int]:
        """슬라이드를 파싱하고 이미지 추출.

        Returns:
            (블록 목록, 이미지 목록, 다음 이미지 인덱스)
        """
        # 모든 요소를 위치 정보와 함께 수집
        elements = []  # [(y_pos, x_pos, element_type, data), ...]
        images = []

        for shape in slide.shapes:
            # 위치 정보 (EMU 단위, None이면 0)
            y_pos = shape.top if shape.top is not None else 0
            x_pos = shape.left if shape.left is not None else 0

            # 이미지 처리
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_block, image_data, image_index = self._extract_picture(
                        shape, image_index
                    )
                    if image_block and image_data:
                        elements.append((y_pos, x_pos, "image", image_block))
                        images.append(image_data)
                except Exception as e:
                    logger.debug("Failed to extract picture from slide %d: %s", slide_idx, e)
            # 테이블 처리
            elif shape.has_table:
                block = self._parse_table(shape.table)
                if block:
                    elements.append((y_pos, x_pos, "table", block))
            # 텍스트 프레임 처리
            elif shape.has_text_frame:
                shape_blocks = self._parse_text_frame(shape, slide_idx)
                for block in shape_blocks:
                    elements.append((y_pos, x_pos, "text", block))

        # y좌표 → x좌표 순으로 정렬 (위에서 아래로, 왼쪽에서 오른쪽으로)
        elements.sort(key=lambda e: (e[0], e[1]))

        # 블록만 추출
        blocks = [elem[3] for elem in elements]

        return blocks, images, image_index

    def _extract_picture(self, shape, image_index: int) -> tuple:
        """Picture shape에서 이미지 추출.

        Returns:
            (ImageBlock, ImageData, 다음 이미지 인덱스)
        """
        try:
            image = shape.image
            image_data = image.blob
            content_type = image.content_type

            ext = get_extension_from_mime(content_type)
            if ext == ".bin":  # 지원하지 않는 이미지 타입
                return None, None, image_index

            image_id = generate_image_id(image_data, image_index)
            image_id_with_ext = f"{image_id}{ext}"

            image_block = create_image(image_id=image_id_with_ext, mime_type=content_type)
            extracted_image = ImageData(
                image_id=image_id_with_ext,
                data=image_data,
                mime_type=content_type,
            )

            return image_block, extracted_image, image_index + 1
        except Exception as e:
            logger.debug("Failed to extract picture data: %s", e)
            return None, None, image_index

    def _parse_text_frame(self, shape, slide_idx: int) -> list:
        """텍스트 프레임을 파싱."""
        blocks = []

        # placeholder 타입 확인
        is_title = False
        is_subtitle = False

        if hasattr(shape, "is_placeholder") and shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            # PP_PLACEHOLDER_TYPE: 1=TITLE, 2=BODY, 3=CENTER_TITLE, 4=SUBTITLE
            if ph_type in [1, 3]:  # TITLE, CENTER_TITLE
                is_title = True
            elif ph_type == 4:  # SUBTITLE
                is_subtitle = True

        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            if is_title:
                # 첫 슬라이드 타이틀은 h1, 나머지는 h2
                level = 1 if slide_idx == 0 else 2
                blocks.append(create_heading(level, text))
            elif is_subtitle:
                blocks.append(create_heading(3, text))
            else:
                # 폰트 크기로 heading 판별 시도
                font_size = self._get_paragraph_font_size(para)
                if font_size and font_size >= Pt(24):
                    blocks.append(create_heading(2, text))
                elif font_size and font_size >= Pt(18):
                    blocks.append(create_heading(3, text))
                else:
                    blocks.append(create_paragraph(text))

        return blocks

    def _get_paragraph_font_size(self, para):
        """문단의 폰트 크기 추출."""
        for run in para.runs:
            if run.font.size:
                return run.font.size
        return None

    def _parse_table(self, table):
        """테이블을 IR 블록으로 변환 (머지셀 포함)."""
        rows_data = []
        merge_info = []
        seen_cells = set()

        for row_idx, row in enumerate(table.rows):
            row_data = []
            col_idx = 0

            for cell in row.cells:
                # 이미 처리된 셀 건너뛰기 (병합된 셀)
                cell_key = (row_idx, col_idx)
                if cell_key in seen_cells:
                    col_idx += 1
                    continue

                cell_text = cell.text.strip() if cell.text else ""
                row_data.append(cell_text)

                # 머지셀 확인
                if cell.is_merge_origin:
                    # 이 셀이 병합의 시작점
                    span_h = cell.span_width  # colspan
                    span_v = cell.span_height  # rowspan

                    if span_h > 1 or span_v > 1:
                        merge_info.append(
                            MergeInfo(
                                row=row_idx,
                                col=len(row_data) - 1,  # 현재 추가된 셀의 인덱스
                                rowspan=span_v,
                                colspan=span_h,
                            )
                        )

                    # 병합된 영역 표시
                    for r in range(row_idx, row_idx + span_v):
                        for c in range(col_idx, col_idx + span_h):
                            seen_cells.add((r, c))

                col_idx += 1

            rows_data.append(row_data)

        if not rows_data:
            return None

        # 첫 행은 header
        header = rows_data[0]
        body_rows = rows_data[1:] if len(rows_data) > 1 else []

        return create_table(
            header=header,
            rows=body_rows,
            merge_info=merge_info if merge_info else None,
        )

    def _extract_charts(
        self, file_path: Path, start_index: int
    ) -> tuple[list, list[ImageData]]:
        """PPTX ZIP 구조에서 차트 추출.

        PPTX 차트 구조:
        - ppt/charts/chart1.xml (차트 데이터)
        - ppt/media/image*.png (차트 이미지)

        Returns:
            (차트 블록 목록, 이미지 데이터 목록)
        """
        blocks = []
        images = []
        image_index = start_index

        try:
            with zipfile.ZipFile(file_path, "r") as zf:
                # 차트 XML 파일 찾기
                chart_files = sorted(
                    [f for f in zf.namelist() if f.startswith("ppt/charts/chart") and f.endswith(".xml")]
                )

                for chart_file in chart_files:
                    try:
                        chart_content = zf.read(chart_file).decode("utf-8")
                        chart_block = parse_chart_xml(chart_content)
                        if chart_block:
                            blocks.append(chart_block)
                    except Exception as e:
                        logger.debug("Failed to parse chart %s: %s", chart_file, e)

                # 차트 이미지 추출 (ppt/media/에서 - 이미 추출된 것 제외)
                # 차트와 연관된 이미지만 추출하기 위해 chart 관련 파일 확인
                for file_name in zf.namelist():
                    if not file_name.startswith("ppt/media/"):
                        continue

                    ext = Path(file_name).suffix.lower()
                    if ext in [".png", ".jpg", ".jpeg", ".gif"]:
                        try:
                            image_data = zf.read(file_name)
                            mime_types = {
                                ".png": "image/png",
                                ".jpg": "image/jpeg",
                                ".jpeg": "image/jpeg",
                                ".gif": "image/gif",
                            }
                            mime_type = mime_types.get(ext, "image/png")

                            image_id = generate_image_id(image_data, image_index)
                            image_id_with_ext = f"{image_id}{ext}"

                            blocks.append(create_image(image_id_with_ext, mime_type))
                            images.append(
                                ImageData(
                                    image_id=image_id_with_ext,
                                    data=image_data,
                                    mime_type=mime_type,
                                )
                            )
                            image_index += 1
                        except Exception as e:
                            logger.debug("Failed to extract media image %s: %s", file_name, e)

        except zipfile.BadZipFile:
            pass

        return blocks, images
