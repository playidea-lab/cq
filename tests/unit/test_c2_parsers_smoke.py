"""Smoke tests for C2 parsers (PDF, DOCX, PPTX).

Minimal "crash 안 함" tests: each parser gets a valid minimal file
and an empty/invalid file. The goal is to prove parsing does not crash,
not to validate content accuracy.
"""

from pathlib import Path

import pytest

# ---------------------------------------------------------------------------
# Helpers: create minimal test files
# ---------------------------------------------------------------------------

def _create_minimal_pdf(path: Path) -> Path:
    """Create a minimal valid PDF with one page containing 'Hello'."""
    from fpdf import FPDF

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    pdf.cell(text="Hello")
    pdf_path = path / "minimal.pdf"
    pdf.output(str(pdf_path))
    return pdf_path


def _create_minimal_docx(path: Path) -> Path:
    """Create a minimal .docx with one paragraph using python-docx."""
    from docx import Document as DocxDocument

    doc = DocxDocument()
    doc.add_paragraph("Hello from DOCX")
    docx_path = path / "minimal.docx"
    doc.save(str(docx_path))
    return docx_path


def _create_minimal_pptx(path: Path) -> Path:
    """Create a minimal .pptx with one slide using python-pptx."""
    from pptx import Presentation

    prs = Presentation()
    slide_layout = prs.slide_layouts[0]  # title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    if title is not None:
        title.text = "Hello from PPTX"
    pptx_path = path / "minimal.pptx"
    prs.save(str(pptx_path))
    return pptx_path


def _create_empty_file(path: Path, name: str) -> Path:
    """Create a zero-byte file."""
    fp = path / name
    fp.write_bytes(b"")
    return fp


# ---------------------------------------------------------------------------
# PDF Parser Smoke Tests
# ---------------------------------------------------------------------------

class TestPdfParserSmoke:
    """Smoke tests for PdfParser."""

    def test_parse_valid_minimal(self, tmp_path: Path) -> None:
        pytest.importorskip("opendataloader_pdf")
        from c4.c2.parsers.pdf_parser import PdfParser

        pdf_path = _create_minimal_pdf(tmp_path)
        parser = PdfParser()
        doc = parser.parse(pdf_path)

        assert doc is not None
        assert isinstance(doc.blocks, list)

    def test_parse_empty_file(self, tmp_path: Path) -> None:
        pytest.importorskip("opendataloader_pdf")
        from c4.c2.parsers.pdf_parser import PdfParser

        empty_path = _create_empty_file(tmp_path, "empty.pdf")
        parser = PdfParser()

        # Empty/invalid file should return empty document (graceful fallback)
        doc = parser.parse(empty_path)
        assert doc is not None
        assert isinstance(doc.blocks, list)


# ---------------------------------------------------------------------------
# DOCX Parser Smoke Tests
# ---------------------------------------------------------------------------

class TestDocxParserSmoke:
    """Smoke tests for DocxParser."""

    def test_parse_valid_minimal(self, tmp_path: Path) -> None:
        pytest.importorskip("docx")
        from c4.c2.parsers.docx_parser import DocxParser

        docx_path = _create_minimal_docx(tmp_path)
        parser = DocxParser()
        doc = parser.parse(docx_path)

        assert doc is not None
        assert isinstance(doc.blocks, list)
        # Should have at least one block from the paragraph we added
        assert len(doc.blocks) >= 1

    def test_parse_empty_file(self, tmp_path: Path) -> None:
        pytest.importorskip("docx")
        from c4.c2.parsers.docx_parser import DocxParser

        empty_path = _create_empty_file(tmp_path, "empty.docx")
        parser = DocxParser()

        # Empty file is not a valid ZIP/DOCX -- should raise cleanly
        with pytest.raises(Exception):
            parser.parse(empty_path)


# ---------------------------------------------------------------------------
# PPTX Parser Smoke Tests
# ---------------------------------------------------------------------------

class TestPptxParserSmoke:
    """Smoke tests for PptxParser."""

    def test_parse_valid_minimal(self, tmp_path: Path) -> None:
        pytest.importorskip("pptx")
        from c4.c2.parsers.pptx_parser import PptxParser

        pptx_path = _create_minimal_pptx(tmp_path)
        parser = PptxParser()
        doc = parser.parse(pptx_path)

        assert doc is not None
        assert isinstance(doc.blocks, list)
        # Should have at least one block from the title we set
        assert len(doc.blocks) >= 1

    def test_parse_empty_file(self, tmp_path: Path) -> None:
        pytest.importorskip("pptx")
        from c4.c2.parsers.pptx_parser import PptxParser

        empty_path = _create_empty_file(tmp_path, "empty.pptx")
        parser = PptxParser()

        # Empty file is not a valid ZIP/PPTX -- should raise cleanly
        with pytest.raises(Exception):
            parser.parse(empty_path)
