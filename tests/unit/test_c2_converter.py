"""Integration tests for C2 converter pipeline (ParserDispatcher).

Tests the full conversion pipeline including:
- Dispatcher routing to correct parsers
- parse() and parse_with_images() for DOCX, PPTX, XLSX
- Error handling for unsupported/missing/corrupt files
"""

from pathlib import Path

import pytest

# ---------------------------------------------------------------------------
# Fixtures: create minimal test files
# ---------------------------------------------------------------------------

@pytest.fixture
def minimal_docx(tmp_path: Path) -> Path:
    """Create a minimal .docx with one paragraph using python-docx."""
    pytest.importorskip("docx")
    from docx import Document as DocxDocument

    doc = DocxDocument()
    doc.add_paragraph("Hello from DOCX")
    doc.add_paragraph("Second paragraph for testing")
    docx_path = tmp_path / "test.docx"
    doc.save(str(docx_path))
    return docx_path


@pytest.fixture
def minimal_pptx(tmp_path: Path) -> Path:
    """Create a minimal .pptx with one slide using python-pptx."""
    pytest.importorskip("pptx")
    from pptx import Presentation

    prs = Presentation()
    slide_layout = prs.slide_layouts[0]  # title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    if title is not None:
        title.text = "Test Presentation"
    pptx_path = tmp_path / "test.pptx"
    prs.save(str(pptx_path))
    return pptx_path


@pytest.fixture
def minimal_xlsx(tmp_path: Path) -> Path:
    """Create a minimal .xlsx with one sheet using openpyxl."""
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    assert ws is not None
    ws.title = "Sheet1"
    ws["A1"] = "Header"
    ws["A2"] = "Data1"
    ws["B1"] = "Value"
    ws["B2"] = 42
    xlsx_path = tmp_path / "test.xlsx"
    wb.save(str(xlsx_path))
    return xlsx_path


@pytest.fixture
def corrupt_file(tmp_path: Path) -> Path:
    """Create a file with invalid content."""
    corrupt_path = tmp_path / "corrupt.docx"
    corrupt_path.write_bytes(b"This is not a valid DOCX file")
    return corrupt_path


# ---------------------------------------------------------------------------
# Test Dispatcher: supported/unsupported extensions
# ---------------------------------------------------------------------------

class TestParserDispatcher:
    """Test ParserDispatcher routing and error handling."""

    def test_unsupported_extension(self, tmp_path: Path) -> None:
        """Dispatcher should raise ValueError for unsupported file types."""
        from c4.c2.parsers.dispatcher import ParserDispatcher

        unsupported_file = tmp_path / "test.xyz"
        unsupported_file.write_text("dummy content")

        dispatcher = ParserDispatcher()
        with pytest.raises(ValueError, match="Unsupported file format"):
            dispatcher.parse(unsupported_file)

    def test_missing_file(self, tmp_path: Path) -> None:
        """Dispatcher should handle missing files gracefully."""
        from c4.c2.parsers.dispatcher import ParserDispatcher

        missing_file = tmp_path / "nonexistent.docx"

        dispatcher = ParserDispatcher()
        with pytest.raises(Exception):  # FileNotFoundError or similar
            dispatcher.parse(missing_file)


# ---------------------------------------------------------------------------
# Test DOCX parsing
# ---------------------------------------------------------------------------

class TestDocxParsing:
    """Test DOCX parsing through Dispatcher."""

    def test_parse_docx(self, minimal_docx: Path) -> None:
        """Dispatcher.parse() should return Document for valid DOCX."""
        pytest.importorskip("docx")
        from c4.c2.parsers.dispatcher import ParserDispatcher

        dispatcher = ParserDispatcher()
        doc = dispatcher.parse(minimal_docx)

        assert doc is not None
        assert hasattr(doc, "blocks")
        assert isinstance(doc.blocks, list)
        assert len(doc.blocks) >= 1

    def test_parse_with_images_docx(self, minimal_docx: Path) -> None:
        """Dispatcher.parse_with_images() should return ParseResult for DOCX."""
        pytest.importorskip("docx")
        from c4.c2.parsers.dispatcher import ParserDispatcher

        dispatcher = ParserDispatcher()
        result = dispatcher.parse_with_images(minimal_docx)

        assert result is not None
        assert hasattr(result, "document")
        assert hasattr(result, "images")
        assert result.document is not None
        assert isinstance(result.images, list)

    def test_corrupt_docx(self, corrupt_file: Path) -> None:
        """Dispatcher should raise exception for corrupt DOCX."""
        pytest.importorskip("docx")
        from c4.c2.parsers.dispatcher import ParserDispatcher

        dispatcher = ParserDispatcher()
        with pytest.raises(Exception):  # BadZipFile or similar
            dispatcher.parse(corrupt_file)


# ---------------------------------------------------------------------------
# Test PPTX parsing
# ---------------------------------------------------------------------------

class TestPptxParsing:
    """Test PPTX parsing through Dispatcher."""

    def test_parse_pptx(self, minimal_pptx: Path) -> None:
        """Dispatcher.parse() should return Document for valid PPTX."""
        pytest.importorskip("pptx")
        from c4.c2.parsers.dispatcher import ParserDispatcher

        dispatcher = ParserDispatcher()
        doc = dispatcher.parse(minimal_pptx)

        assert doc is not None
        assert hasattr(doc, "blocks")
        assert isinstance(doc.blocks, list)
        assert len(doc.blocks) >= 1

    def test_parse_with_images_pptx(self, minimal_pptx: Path) -> None:
        """Dispatcher.parse_with_images() should return ParseResult for PPTX."""
        pytest.importorskip("pptx")
        from c4.c2.parsers.dispatcher import ParserDispatcher

        dispatcher = ParserDispatcher()
        result = dispatcher.parse_with_images(minimal_pptx)

        assert result is not None
        assert hasattr(result, "document")
        assert hasattr(result, "images")
        assert result.document is not None
        assert isinstance(result.images, list)


# ---------------------------------------------------------------------------
# Test XLSX parsing
# ---------------------------------------------------------------------------

class TestXlsxParsing:
    """Test XLSX parsing through Dispatcher."""

    def test_parse_xlsx(self, minimal_xlsx: Path) -> None:
        """Dispatcher.parse() should return Document for valid XLSX."""
        pytest.importorskip("openpyxl")
        from c4.c2.parsers.dispatcher import ParserDispatcher

        dispatcher = ParserDispatcher()
        doc = dispatcher.parse(minimal_xlsx)

        assert doc is not None
        assert hasattr(doc, "blocks")
        assert isinstance(doc.blocks, list)
        # XLSX should have at least table blocks
        assert len(doc.blocks) >= 1

    def test_parse_with_images_xlsx(self, minimal_xlsx: Path) -> None:
        """Dispatcher.parse_with_images() should return ParseResult for XLSX."""
        pytest.importorskip("openpyxl")
        from c4.c2.parsers.dispatcher import ParserDispatcher

        dispatcher = ParserDispatcher()
        result = dispatcher.parse_with_images(minimal_xlsx)

        assert result is not None
        assert hasattr(result, "document")
        assert hasattr(result, "images")
        assert result.document is not None
        assert isinstance(result.images, list)
